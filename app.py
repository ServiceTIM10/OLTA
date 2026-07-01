import re
import tempfile
import unicodedata
from pathlib import Path
from io import BytesIO
from email import policy
from email.parser import BytesParser
from email.utils import parsedate_to_datetime

import pandas as pd
import streamlit as st
from bs4 import BeautifulSoup

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
    CategoryChartData = None
    RGBColor = None
    PP_ALIGN = None
    MSO_ANCHOR = None
    Inches = None
    Pt = None

from copy import deepcopy

try:
    import extract_msg
except Exception:
    extract_msg = None

FINAL_COLUMNS = [
    "OLTA MITTENTE",
    "DATA NEWSLETTER",
    "TIPOLOGIA",
    "COMPAGNIA PROMOSSA",
    "DESTINAZIONE/MERCATO/TRATTA",
    "PROMOZIONE",
    "OGGETTO EMAIL",
    "FILE ORIGINE",
    "TESTO ESTRATTO",
    "CONFIDENCE",
    "NOTE",
]

# Colonne visibili nelle sotto-tabelle dell'app.
# Le colonne tecniche FILE ORIGINE, CONFIDENCE, TESTO ESTRATTO e NOTE restano disponibili
# internamente, ma non vengono mostrate come colonne nelle tabelle modificabili.
DISPLAY_COLUMNS = [
    "COMPAGNIA PROMOSSA",
    "DESTINAZIONE/MERCATO/TRATTA",
    "PROMOZIONE",
    "OGGETTO EMAIL",
]

# Colonne esportate nel file Excel finale.
# TESTO ESTRATTO e NOTE sono escluse perché richiesto di eliminarle dalle tabelle operative.
EXPORT_COLUMNS = [
    "OLTA MITTENTE",
    "DATA NEWSLETTER",
    "TIPOLOGIA",
    "COMPAGNIA PROMOSSA",
    "DESTINAZIONE/MERCATO/TRATTA",
    "PROMOZIONE",
    "OGGETTO EMAIL",
    "FILE ORIGINE",
    "CONFIDENCE",
]

KNOWN_OLTAS = {
    "AFERRY": ["aferry", "anna aferry"],
    "ALLOFERRY": ["allo ferry", "alloferry"],
    "FERRYHOPPER": ["ferryhopper"],
    "NETFERRY": ["netferry"],
    "TRAGHETTIPER": ["traghettiper"],
    "TRAGHETTI.COM": ["traghetti.com"],
    "DIRECT FERRIES": ["direct ferries", "directferries", "direct_ferries"],
    "LA CENTRALE DES FERRIES": ["la centrale des ferries", "lacentrale", "la centrale", "lacentraledesferries"],
}

APP_DIR = Path(__file__).parent
LOGO_WIDTH_PX = 165

# Streamlit Community Cloud usa Linux: i nomi file sono case-sensitive.
# Per evitare problemi tra "assets/logos" e "logos", maiuscole/minuscole o spazi,
# l'app cerca i loghi in più percorsi e con più varianti di nome.
LOGO_DIR_CANDIDATES = [
    APP_DIR / "assets" / "logos",
    APP_DIR / "logos",
    Path.cwd() / "assets" / "logos",
    Path.cwd() / "logos",
]

OLTA_LOGO_FILES = {
    "AFERRY": ["aferry.png", "AFERRY.png"],
    "ALLOFERRY": ["alloferry.png", "ALLOFERRY.png", "allo_ferry.png", "ALLO FERRY.png"],
    "FERRYHOPPER": ["ferryhopper.png", "FERRYHOPPER.png", "ferry_hopper.png"],
    "NETFERRY": ["netferry.png", "NETFERRY.png"],
    "TRAGHETTIPER": ["traghettiper.png", "TRAGHETTIPER.png", "traghetti_per.png"],
    "TRAGHETTI.COM": ["traghetti_com.png", "TRAGHETTICOM.png", "TRAGHETTI.COM.png", "traghetti.com.png"],
    "DIRECT FERRIES": ["direct_ferries.png", "DIRECT FERRIES.png", "direct ferries.png"],
    "LA CENTRALE DES FERRIES": ["la_centrale_des_ferries.png", "LA CENTRALE DES FERRIES.png"],
}

TEMPLATE_PPTX_CANDIDATES = [
    APP_DIR / "templates" / "template_olta.pptx",
    APP_DIR / "template_olta.pptx",
    Path.cwd() / "templates" / "template_olta.pptx",
    Path.cwd() / "template_olta.pptx",
]

PPT_OLTA_CATEGORIES = [
    "ADAC",
    "AFERRY",
    "ALLOFERRY",
    "DIRECT FERRIES",
    "FERRY HOPPER",
    "LA CENTRALE DES FERRIES",
    "NETFERRY",
    "TRAGHETTIPER",
    "TRAGHETTI.COM",
]

PPT_BAR_COLORS = [
    (31, 78, 121),
    (112, 173, 71),
    (237, 125, 49),
    (91, 155, 213),
    (165, 165, 165),
    (255, 192, 0),
    (68, 114, 196),
    (112, 48, 160),
    (0, 176, 80),
]

KNOWN_COMPANIES = [
    "GNV", "Grimaldi Lines", "Moby", "Tirrenia", "Corsica Ferries", "Sardinia Ferries",
    "Corsica Linea", "La Méridionale", "La Meridionale", "Tallink Silja Line",
    "Trasmed GLE", "Trasmed", "Brittany Ferries", "DFDS", "Irish Ferries", "P&O Ferries",
    "P&O", "Color Line", "Superfast Ferries", "Liberty Lines", "Siremar", "Stena Line",
    "Armas", "Balearia", "SNAV", "Jadrolinija", "Blue Star Ferries", "Seajets",
    "Minoan Lines", "ANEK Lines", "Viking Line", "Finnlines",
    "Lafasi", "Ventouris Ferries", "Anek Superfast",
    "Magic Sea Ferries", "Golden Star Ferries", "Tallink Silja",
]

KNOWN_DESTINATIONS = [
    "Sicilia", "Sardegna", "Grecia", "Italia-Grecia", "Italia ↔ Grecia", "Italia Grecia",
    "Inghilterra", "Irlanda", "Baleari", "Isole Baleari", "Ibiza", "Maiorca", "Minorca",
    "Tunisia", "Algeria", "Marocco", "Francia-Marocco", "Francia - Marocco",
    "Francia Algérie", "France Algérie", "Danimarca", "Norvegia", "Danimarca ↔ Norvegia",
    "Mar Baltico", "Corsica", "Elba", "Isola d'Elba", "Isole Pelagie", "Lampedusa", "Linosa",
    "Porto Empedocle", "Nord Europa", "Mare d'Irlanda", "Manica", "Fiordi",
    "Almeria Nador", "Motril Nador", "Almeria-Melilla", "Motril-Melilla", "Almeria Oran",
    "Almeria Ghazouet", "Algéciras Tanger", "Algésiras Tanger", "Algéciras Ceuta",
    "Sète Tanger", "Sète Nador", "Barcelone Tanger", "Barcelone Nador", "Gênes Tanger",
    "Marseille", "Sete", "Sète", "Alger", "Bejaia", "Nador", "Tanger", "Ceuta", "Oran", "Ghazouet",
    "Albania", "Croazia", "Croazia e isole italiane", "Isole Greche", "Isole Cicladi",
    "Svezia", "Finlandia", "Estonia", "Grecia - Turchia", "Grecia-Turchia", "Marocco", "Maroc",
]

COMMON_LINES = {
    "Prenota ora", "Trova il miglior prezzo", "Cerca la tua tratta su TraghettiPer.",
    "Confronta le offerte di tutte le compagnie di navigazione.", "Buona navigazione,",
    "Il Team di TraghettiPer", "Perché viaggiare con noi", "Professionale", "Sicuro", "Veloce",
    "Dicono di noi", "Metodi di pagamento", "Annulla iscrizione", "Cancella iscrizione alla newsletter",
}

PROMO_RE = re.compile(
    r"(?i)(?:fino\s+al|fino\s+a|jusqu['’]à|up\s+to|sconto\s+del|sconto|remise|"
    r"risparmia(?:\s+fino\s+al|\s+il)?|offerta|promo|coupon|codice|voucher|bon\s+d[’']achat|"
    r"a\s+partire\s+da|da|dès|à\s+partir\s+de|partire\s+da|-\s?\d{1,2}\s?%|"
    r"\d{1,2}\s?%|\d+(?:[,.]\d+)?\s?€|veicolo\s+gratis|gratis)"
)
COUPON_CODE_RE = re.compile(r"(?i)(?:codice|coupon|voucher|bon\s+d[’']achat)")
ROUTE_HINT_RE = re.compile(r"(?i)(?:\b[A-ZÀ-Ý][a-zà-ÿ]+(?:\s*[↔\-–]\s*[A-ZÀ-Ý][a-zà-ÿ]+)+\b)")


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = str(text)
    text = re.sub(r"<[^>]*(?:https?|mailto):[^>]+>", " ", text)
    text = re.sub(r"https?://\S+", " ", text)
    text = re.sub(r"mailto:\S+", " ", text)
    text = re.sub(r"[\u00ad\u034f\u061c\u180e\u200b-\u200f\u202a-\u202e\u2060-\u206f\ufeff]", " ", text)
    text = text.replace("\xa0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n", text)
    return text.strip()


def raw_clean_lines(text: str) -> list[str]:
    text = clean_text(text)
    raw = [line.strip(" \t\r\n-*•·") for line in text.splitlines()]
    lines = []
    for line in raw:
        line = re.sub(r"\s+", " ", line).strip()
        if not line or len(line) < 2:
            continue
        if line.lower().startswith("http"):
            continue
        if re.search(
            r"(?i)annulla iscrizione|unsubscribe|copyright|tutti i diritti|ricevi questa email|"
            r"vous recevez ce mail|conformément|dicono di noi|metodi di pagamento|facebook|instagram|you received",
            line,
        ):
            continue
        lines.append(line)
    return lines


def filtered_lines(text: str) -> list[str]:
    return [line for line in raw_clean_lines(text) if line not in COMMON_LINES]


def html_to_text_with_image_metadata(html: str) -> str:
    """Estrae il testo HTML aggiungendo anche alt/title delle immagini.

    Alcune newsletter hanno una parte rilevante dentro immagini; quando alt/title
    sono disponibili, li aggiungiamo al testo per facilitare le regole.
    """
    if not html:
        return ""
    soup = BeautifulSoup(html, "html.parser")
    text_parts = [soup.get_text("\n", strip=True)]
    image_texts = []
    for img in soup.find_all("img"):
        for attr in ["alt", "title"]:
            value = img.get(attr)
            if value:
                value = cleanup_field(value)
                if value and value not in image_texts:
                    image_texts.append(value)
    if image_texts:
        text_parts.append("\n".join(image_texts))
    return clean_text("\n".join(part for part in text_parts if part))


def parse_uploaded_file(uploaded_file) -> dict:
    filename = uploaded_file.name
    suffix = Path(filename).suffix.lower()
    data = uploaded_file.getvalue()

    if suffix == ".msg":
        if extract_msg is None:
            raise RuntimeError("Per leggere i file .msg devi installare la libreria extract-msg.")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tmp:
            tmp.write(data)
            tmp_path = tmp.name
        msg = extract_msg.Message(tmp_path)
        html_body = msg.htmlBody
        html_text = ""
        if isinstance(html_body, bytes):
            html_text = html_body.decode("utf-8", errors="ignore")
        elif html_body:
            html_text = str(html_body)
        body = msg.body or ""
        if body and len(body.strip()) > 30:
            # Se il body testuale è disponibile, lo usiamo come fonte principale.
            # Gli alt/title delle immagini vengono usati solo quando manca il body,
            # così evitiamo duplicazioni/false offerte nei template ricchi di immagini.
            text = body
        elif html_text:
            text = html_to_text_with_image_metadata(html_text)
        else:
            text = ""
        return {"sender": msg.sender or "", "date": msg.date, "subject": msg.subject or "", "text": clean_text(text)}

    if suffix == ".eml":
        msg = BytesParser(policy=policy.default).parsebytes(data)
        sender = msg.get("from", "")
        subject = msg.get("subject", "")
        try:
            date = parsedate_to_datetime(msg.get("date")) if msg.get("date") else None
        except Exception:
            date = None
        plain_parts, html_parts = [], []
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_disposition() == "attachment":
                    continue
                try:
                    content = part.get_content()
                except Exception:
                    continue
                if part.get_content_type() == "text/plain":
                    plain_parts.append(content)
                elif part.get_content_type() == "text/html":
                    html_parts.append(content)
        else:
            content = msg.get_content()
            if msg.get_content_type() == "text/plain":
                plain_parts.append(content)
            elif msg.get_content_type() == "text/html":
                html_parts.append(content)
        text = "\n".join(plain_parts) if plain_parts else "\n".join(
            html_to_text_with_image_metadata(h) for h in html_parts
        )
        return {"sender": sender, "date": date, "subject": subject, "text": clean_text(text)}

    if suffix in [".html", ".htm"]:
        raw = data.decode("utf-8", errors="ignore")
        text = html_to_text_with_image_metadata(raw)
        return {"sender": "", "date": None, "subject": "", "text": clean_text(text)}

    if suffix == ".txt":
        raw = data.decode("utf-8", errors="ignore")
        return {"sender": "", "date": None, "subject": "", "text": clean_text(raw)}

    raw = data.decode("utf-8", errors="ignore")
    return {"sender": "", "date": None, "subject": "", "text": clean_text(raw)}


def get_olta(sender: str, filename: str) -> str:
    source = f"{sender} {filename}".lower()
    for olta, aliases in KNOWN_OLTAS.items():
        if any(alias.lower() in source for alias in aliases):
            return olta
    stem = Path(filename).stem
    match = re.match(r"([A-Za-z ._-]+)[_\-\s]\d", stem)
    if match:
        return match.group(1).replace("_", " ").strip().upper()
    return ""


def format_date(dt, filename: str) -> str:
    if dt:
        try:
            return dt.strftime("%d/%m/%Y")
        except Exception:
            pass
    match = re.search(r"(\d{1,2})[.\-_](\d{1,2})[.\-_](\d{4})", filename)
    if match:
        return f"{int(match.group(1)):02d}/{int(match.group(2)):02d}/{match.group(3)}"
    match = re.search(r"(\d{4})[.\-_](\d{1,2})[.\-_](\d{1,2})", filename)
    if match:
        return f"{int(match.group(3)):02d}/{int(match.group(2)):02d}/{match.group(1)}"
    return ""


def cleanup_field(value) -> str:
    if not value:
        return ""
    value = str(value)
    value = re.sub(r"\s+", " ", value).strip(" -;:,")
    return value.replace(" ,", ",")


def find_companies(context: str) -> str:
    found = []
    for company in KNOWN_COMPANIES:
        pattern = re.compile(r"(?i)(?<![A-Za-z])" + re.escape(company).replace("\\ ", r"\s+") + r"(?![A-Za-z])")
        if pattern.search(context):
            normalized = company.replace("La Meridionale", "La Méridionale")
            found.append(normalized)
    output = []
    for company in found:
        if company == "P&O" and "P&O Ferries" in found:
            continue
        if company not in output:
            output.append(company)
    return ", ".join(output)


def find_destinations(context: str) -> str:
    found = []
    for dest in KNOWN_DESTINATIONS:
        pattern = re.compile(r"(?i)(?<![A-Za-zÀ-ÿ])" + re.escape(dest).replace("\\ ", r"\s+") + r"(?![A-Za-zÀ-ÿ])")
        if pattern.search(context):
            found.append(dest)
    for route in ROUTE_HINT_RE.findall(context):
        if route not in found:
            found.append(route)
    return ", ".join(dict.fromkeys(found))


def extract_promo_phrase(context: str) -> str:
    context = " ".join(line.strip() for line in str(context).splitlines() if line.strip())

    match = re.search(r"(?i)(?:valore\s+del\s+coupon\s*:?\s*)?(\d+(?:[,.]\d+)?\s?€\s*(?:di\s+)?sconto(?:\s*\([^)]*\))?)", context)
    if match:
        return match.group(1).strip()

    if re.search(r"(?i)meilleur prix garanti|miglior prezzo garantito|best price guarantee", context):
        return "Meilleur Prix Garanti / rimborso differenza"

    match = re.search(r"(?i)(biglietto\s+da\s+\d+(?:[,.]\d+)?\s?€\s*\+\s*veicolo\s+gratis(?:\s+al\s+ritorno)?)", context)
    if match:
        return match.group(1).strip()

    match = re.search(r"(?i)(veicolo\s+gratis(?:\s+al\s+ritorno)?|vehicle\s+free|véhicule\s+gratuit)", context)
    if match:
        return match.group(1).strip()

    patterns = [
        r"(?i)(fino\s+al\s+\d{1,2}\s?%\s*(?:di\s+sconto)?)",
        r"(?i)(fino\s+a\s+\d{1,2}\s?%\s*(?:di\s+sconto)?)",
        r"(?i)(risparmia\s+fino\s+al\s+\d{1,2}\s?%)",
        r"(?i)(risparmia\s+il\s+\d{1,2}\s?%)",
        r"(?i)(sconto\s+del\s+\d{1,2}\s?%)",
        r"(?i)(\d{1,2}\s?%\s+di\s+sconto)",
        r"(?i)(-\s?\d{1,2}\s?%)",
        r"(?i)(jusqu['’]à\s+\d{1,2}\s?%)",
        r"(?i)(remise\s+.*?\d{1,2}\s?%)",
    ]
    for pattern in patterns:
        match = re.search(pattern, context)
        if match:
            return re.sub(r"\s+", " ", match.group(1)).strip()

    match = re.search(r"(?i)((?:a\s+partire\s+da|partire\s+da|da|dès|à\s+partir\s+de)\s*\d+(?:[,.]\d+)?\s?€)", context)
    if match:
        return match.group(1).strip()

    return ""


def infer_tipologia(subject: str, text: str, promo: str) -> str:
    sample = f"{subject}\n{text[:4000]}".lower()
    if re.search(r"\b(coupon|codice|voucher|bon d’achat|code promo|coupon sconto)\b", sample):
        return "Codice sconto"
    if re.search(r"\b(ultimi giorni|ultimo giorno|last chance|solo oggi|demain|derni[eè]re ligne droite)\b", sample):
        return "Reminder"
    if promo or re.search(r"\b(sconto|offerta|promo|risparmia|remise|offre|a partire da|fino al|jusqu|discount)\b|-\s?\d{1,2}\s?%|\d+\s?€", sample):
        return "Promozione"
    return "Comunicazione generica"


def make_row(parsed: dict, filename: str, tipologia: str, company: str, destination: str, promo: str, confidence: float, note: str) -> dict:
    text = parsed.get("text", "") or ""
    return {
        "OLTA MITTENTE": get_olta(parsed.get("sender", ""), filename),
        "DATA NEWSLETTER": format_date(parsed.get("date"), filename),
        "TIPOLOGIA": cleanup_field(tipologia),
        "COMPAGNIA PROMOSSA": cleanup_field(company),
        "DESTINAZIONE/MERCATO/TRATTA": cleanup_field(destination),
        "PROMOZIONE": cleanup_field(promo),
        "OGGETTO EMAIL": parsed.get("subject", "") or "",
        "FILE ORIGINE": filename,
        "TESTO ESTRATTO": (text[:1000] + "…") if len(text) > 1000 else text,
        "CONFIDENCE": round(float(confidence), 2),
        "NOTE": note,
    }


def extract_fallback(parsed: dict, filename: str) -> dict:
    full = f"{parsed.get('subject', '')}\n{parsed.get('text', '')[:4000]}"
    promo = extract_promo_phrase(full) or "Comunicazione / promozione senza meccanica economica esplicita"
    company = find_companies(full)
    destination = find_destinations(full)
    tipologia = infer_tipologia(parsed.get("subject", ""), parsed.get("text", ""), promo)
    return make_row(parsed, filename, tipologia, company, destination, promo, 0.45, "Fallback: estrazione generica")


def extract_aferry(parsed: dict, filename: str) -> list[dict]:
    lines = raw_clean_lines(parsed.get("text", ""))
    positions = [i for i, line in enumerate(lines) if re.search(r"(?i)trova il miglior prezzo", line)]
    start = positions[-1] + 1 if positions else 0
    end = len(lines)
    for i, line in enumerate(lines[start:], start=start):
        if re.search(r"(?i)prenota più veloce|termini e condizioni", line):
            end = i
            break
    section = lines[start:end]
    blocks, current = [], []
    for line in section:
        if re.search(r"(?i)^prenota ora\b", line):
            if current:
                blocks.append(current)
                current = []
        else:
            current.append(line)
    if current:
        blocks.append(current)

    rows, seen = [], set()
    for block in blocks:
        block = [
            line for line in block
            if not re.search(r"(?i)^destinazioni|^traversate|^verso |^le tue prossime|^in famiglia|^viaggia a modo|^verso sud|^è arrivato|^più luce|^scopri le offerte|^baleari, tunisia|^corsica, marocco|^vacanze in famiglia", line)
        ]
        if not block:
            continue
        block_text = "\n".join(block)
        promo = extract_promo_phrase(block_text)
        if not promo:
            continue
        headline = ""
        for line in block:
            if re.search(r"(?i)^prenota ora|^trova|^scopri|^tutte le offerte|^approfitta|^risparmia|^fino al|^sconto|^corsica linea|^tallink|^gnv:|^brittany ferries|^dfds", line):
                continue
            if (":" in line and (PROMO_RE.search(line) or find_destinations(line))) or find_destinations(line):
                headline = line
                break
        if not headline:
            headline = block[0]
        match = re.match(r"^(.+?)\s*:\s*(.+)$", headline)
        destination = match.group(1).strip() if match else headline
        if re.search(r"(?i)tutte le offerte|risparmia|sconto|fino|approfitta", destination):
            destination = find_destinations(block_text)
        company = find_companies(block_text)
        tipologia = infer_tipologia(parsed.get("subject", ""), block_text, promo)
        key = (company, destination, promo)
        if key in seen:
            continue
        seen.add(key)
        rows.append(make_row(parsed, filename, tipologia, company, destination, promo, 0.90, "Regola specifica AFerry: blocco offerta"))
    return rows


def extract_alloferry(parsed: dict, filename: str) -> list[dict]:
    lines = filtered_lines(parsed.get("text", ""))
    full = "\n".join(lines[:120])
    rows = []

    if re.search(r"(?i)meilleur prix garanti", full):
        company = find_companies(full)
        match = re.search(
            r"(?i)Offre valable pour les traversées\s*:\s*(.+?)(?:\n\s*- Offre applicable|\n\s*- La garantie|\n\s*- L'offre)",
            full,
            re.S,
        )
        destination = re.sub(r"\s+", " ", match.group(1)).strip(" -") if match else find_destinations(full)
        rows.append(make_row(parsed, filename, "Promozione", company, destination, "Meilleur Prix Garanti / rimborso differenza", 0.80, "Regola specifica Alloferry: garanzia prezzo"))
        return rows

    company = ""
    match_company = re.search(r"(?i)compagnie\s*:\s*([^\n]+)", full)
    if match_company:
        company = match_company.group(1).strip()
    else:
        company = find_companies(full)

    seen = set()
    for line in lines[:40]:
        if not re.search(r"\d{1,2}\s?%", line):
            continue
        if re.search(r"(?i)prix hors taxe|non rétroactive|non cumulable|valable|places limitées|remise déduite|réduction", line):
            continue
        promo = extract_promo_phrase(line)
        if not promo:
            match_percent = re.search(r"(\d{1,2}\s?%)", line)
            promo = match_percent.group(1) if match_percent else ""
        destination = ""
        match = re.match(r"(?i)^(.+?)\s*[-: ]+\s*\d{1,2}\s?%\s*: ?", line)
        if match and not re.match(r"^\s*\d", line):
            destination = match.group(1).strip(" -:")
        if not destination:
            match = re.search(r"(?i)-?\s*\d{1,2}\s?%\s+(?:sur|sui|su|on)?\s*(.+)$", line)
            if match:
                destination = match.group(1).strip(" -:")
        key = (destination, promo, company)
        if promo and destination and key not in seen:
            rows.append(make_row(parsed, filename, "Promozione", company, destination, promo, 0.85, "Regola specifica Alloferry: headline con percentuale"))
            seen.add(key)
    return rows


def extract_traghettiper(parsed: dict, filename: str) -> list[dict]:
    lines = filtered_lines(parsed.get("text", ""))
    full = "\n".join(lines[:100])
    promo = extract_promo_phrase(full)
    match = re.search(r"(?i)valore\s+del\s+coupon\s*:?\s*([^\n]+)", full)
    if match:
        promo = match.group(1).strip()
    code = ""
    match_code = re.search(r"(?i)(?:il tuo codice traghettiper|codice traghettiper)\s*:?\s*\n?\s*([A-Z0-9-]{4,})", full)
    if match_code:
        code = match_code.group(1).strip()
    destination = find_destinations(f"{parsed.get('subject', '')}\n{full}")
    if not destination and re.search(r"(?i)tutte le prenotazioni", full):
        destination = "Tutte le prenotazioni"
    note = "Regola specifica TraghettiPer: coupon"
    if code:
        note += f" | Codice rilevato: {code}"
    tipologia = "Codice sconto" if COUPON_CODE_RE.search(full) else infer_tipologia(parsed.get("subject", ""), full, promo)
    return [make_row(parsed, filename, tipologia, "", destination, promo or "Coupon/codice sconto", 0.90 if promo else 0.65, note)]


def extract_netferry(parsed: dict, filename: str) -> list[dict]:
    full = f"{parsed.get('subject', '')}\n{parsed.get('text', '')[:3000]}"
    promo = extract_promo_phrase(full)
    company = find_companies(full)
    destination = find_destinations(full)
    if not promo:
        promo = "Comunicazione / promozione senza meccanica economica esplicita"
        tipologia = "Comunicazione generica"
        confidence = 0.50
    else:
        tipologia = infer_tipologia(parsed.get("subject", ""), parsed.get("text", ""), promo)
        confidence = 0.75
    return [make_row(parsed, filename, tipologia, company, destination, promo, confidence, "Regola specifica Netferry: sintesi newsletter")]


def extract_ferryhopper(parsed: dict, filename: str) -> list[dict]:
    lines = filtered_lines(parsed.get("text", ""))
    rows, seen = [], set()
    for i, line in enumerate(lines[:90]):
        if not PROMO_RE.search(line):
            continue
        if re.search(r"(?i)^promo traghetti|fino al -20% sui traghetti", line):
            continue
        if re.search(r"(?i)scade|assistenza|app di ferryhopper", line):
            continue
        prevs = [lines[j] for j in range(max(0, i - 3), i)]
        nexts = [lines[j] for j in range(i + 1, min(len(lines), i + 4))]
        context = "\n".join(prevs + [line] + nexts)
        promo_context = line + "\n" + ("\n".join(nexts[:1]) if re.search(r"(?i)veicolo|gratis", line) else "")
        promo = extract_promo_phrase(promo_context)
        if not promo:
            continue
        destination = ""
        for previous in reversed(prevs):
            if len(previous) > 3 and not PROMO_RE.search(previous) and not re.search(
                r"(?i)super promo|vamos|scopri|prenota|sugli itinerari|valido|traghetti scontati|mix perfetto|calette|segreto|itinerario",
                previous,
            ):
                destination = previous
                break
        if not destination:
            destination = find_destinations(context)
        company = find_companies(context)
        local_company = find_companies("\n".join([line] + nexts))
        if local_company:
            company = local_company
        key = (destination, promo, company)
        if key in seen:
            continue
        seen.add(key)
        rows.append(make_row(parsed, filename, "Promozione", company, destination, promo, 0.85, "Regola specifica Ferryhopper: riga promo + contesto"))
    return rows


def extract_direct_ferries(parsed: dict, filename: str) -> list[dict]:
    """Estrazione specifica per newsletter Direct Ferries.

    Le newsletter Direct Ferries sono organizzate in blocchi separati da CTA
    "Cerca Ora": destinazione, meccanica promozionale e compagnia.
    """
    lines = raw_clean_lines(parsed.get("text", ""))
    cleaned_lines = []
    for line in lines:
        if re.search(r"(?i)caution:|questa mail proviene|direct ferries ltd|cancella l.iscrizione|tutti i diritti|versione online|contattarci", line):
            continue
        if re.search(r"(?i)^prenota oggi|^mare, sole|godersi le vacanze|^fino al 50% di$|^sconto sui traghetti|^sali a bordo|^esplora il nostro blog|^alla scoperta|^spiagge segrete|^per saperne|^nota bene", line):
            continue
        if line.startswith("<data:image"):
            continue
        cleaned_lines.append(line)

    # Split in blocchi tra una CTA e l'altra.
    blocks, current = [], []
    for line in cleaned_lines:
        if re.search(r"(?i)^cerca\s+ora$|^cerca$|^ora$", line):
            if current:
                blocks.append(current)
                current = []
        else:
            current.append(line)
    if current:
        blocks.append(current)

    rows, seen = [], set()
    for block in blocks:
        block = [cleanup_field(x).replace("*", "") for x in block if cleanup_field(x)]
        if not block:
            continue
        block_text = "\n".join(block)
        if not re.search(r"(?i)\bcon\s+[A-Za-zÀ-ÿ]|\d{1,2}\s?%", block_text):
            continue

        # Prima riga realmente descrittiva = destinazione/mercato/tratta.
        destination = ""
        for line in block:
            if re.search(r"(?i)^fino al|^\d{1,2}\s?%|^con\s+|sconto|biglietti", line):
                continue
            if len(line) > 2:
                destination = line.strip(" .")
                break
        if not destination:
            destination = find_destinations(block_text)

        # Caso speciale: blocco con più offerte nella stessa destinazione (es. Grecia).
        bullet_matches = re.findall(r"(?i)(\d{1,2}\s?%\s+di\s+sconto)\s+con\s+([^\n]+)", block_text) if re.search(r"(?im)^Fino\s+al:?$", block_text) else []
        if bullet_matches:
            for promo_part, company_part in bullet_matches:
                promo = promo_part.strip()
                if not re.search(r"(?i)^fino", promo):
                    promo = f"Fino al {promo}"
                company = cleanup_field(company_part)
                key = (destination, company, promo)
                if key not in seen:
                    seen.add(key)
                    rows.append(make_row(parsed, filename, "Promozione", company, destination, promo, 0.95, "Regola specifica Direct Ferries: blocco multi-offerta"))
            continue

        promo = extract_promo_phrase(block_text)
        # Integra condizioni testuali immediatamente successive alla percentuale.
        if promo and re.search(r"(?i)sui biglietti con rientro in giornata", block_text) and "rientro in giornata" not in promo.lower():
            promo = f"{promo} sui biglietti con rientro in giornata"

        company = ""
        match_company = re.search(r"(?im)^con\s+([^\n]+)$", block_text)
        if match_company:
            company = cleanup_field(match_company.group(1))
        if not company:
            company = find_companies(block_text)
        if not promo or not company:
            continue
        key = (destination, company, promo)
        if key in seen:
            continue
        seen.add(key)
        rows.append(make_row(parsed, filename, "Promozione", company, destination, promo, 0.95, "Regola specifica Direct Ferries: blocco offerta"))

    return rows


def extract_la_centrale_des_ferries(parsed: dict, filename: str) -> list[dict]:
    """Estrazione specifica per La Centrale des Ferries.

    Questi template sono spesso image-based: in assenza di testo completo nelle
    immagini, la regola usa oggetto + preheader/testo visibile.
    """
    subject = cleanup_field(parsed.get("subject", ""))
    lines = raw_clean_lines(parsed.get("text", ""))
    useful_lines = []
    for line in lines:
        line_clean = cleanup_field(line)
        if not line_clean or line_clean == subject:
            continue
        if len(line_clean) > 180:
            continue
        if re.search(r"(?i)désinscrire|newsletter|sendibm|contacter|conditions|cliquez ici", line_clean):
            continue
        if re.match(r"^<.*>$", line_clean):
            continue
        useful_lines.append(line_clean)
        if len(useful_lines) >= 5:
            break

    full = cleanup_field("\n".join([subject] + useful_lines))
    full_lower = full.lower()

    destination = find_destinations(full)
    if re.search(r"(?i)\bmaroc\b|tanger|nador", full):
        routes = []
        for route in ["Algeciras-Tanger", "Almeria-Nador", "Almería-Nador", "Sète-Nador", "Sète-Tanger", "Marseille-Tanger"]:
            if re.search(re.escape(route), full, re.I):
                routes.append(route)
        destination = "Marocco" + (" / " + ", ".join(routes) if routes else "")
    elif re.search(r"(?i)alg[eé]rie|alger|oran|bejaia|béjaïa", full):
        destination = "Algeria"

    if re.search(r"(?i)payez en plusieurs fois|paiement en plusieurs fois|ancv", full):
        promo = "pagamento in più rate / ANCV"
    elif re.search(r"(?i)meilleur prix", full):
        promo = "ferry al miglior prezzo"
    elif re.search(r"(?i)promo", full):
        promo = "prenotazione ferry in promo"
    else:
        promo = "comunicazione commerciale / vacanze"

    if re.search(r"(?i)derni[eè]re ligne droite", full):
        tipologia = "Reminder"
    elif re.search(r"(?i)promo|meilleur prix|prix|ancv|payez", full):
        tipologia = "Promozione"
    else:
        tipologia = "Comunicazione generica"

    note = "Regola specifica La Centrale des Ferries: estrazione da oggetto/preheader; template prevalentemente image-based"
    confidence = 0.78 if useful_lines else 0.62
    return [make_row(parsed, filename, tipologia, "", destination, promo, confidence, note)]


def extract_rows(parsed: dict, filename: str) -> list[dict]:
    olta = get_olta(parsed.get("sender", ""), filename)
    try:
        if olta == "AFERRY":
            rows = extract_aferry(parsed, filename)
        elif olta == "ALLOFERRY":
            rows = extract_alloferry(parsed, filename)
        elif olta == "TRAGHETTIPER":
            rows = extract_traghettiper(parsed, filename)
        elif olta == "NETFERRY":
            rows = extract_netferry(parsed, filename)
        elif olta == "FERRYHOPPER":
            rows = extract_ferryhopper(parsed, filename)
        elif olta == "DIRECT FERRIES":
            rows = extract_direct_ferries(parsed, filename)
        elif olta == "LA CENTRALE DES FERRIES":
            rows = extract_la_centrale_des_ferries(parsed, filename)
        else:
            rows = []
        if not rows:
            rows = [extract_fallback(parsed, filename)]
    except Exception as exc:
        row = extract_fallback(parsed, filename)
        row["NOTE"] += f" | Errore regola specifica: {exc}"
        row["CONFIDENCE"] = 0.30
        rows = [row]

    cleaned = []
    seen = set()
    for row in rows:
        for column in FINAL_COLUMNS:
            row.setdefault(column, "")
        key = tuple(row[col] for col in ["OLTA MITTENTE", "DATA NEWSLETTER", "COMPAGNIA PROMOSSA", "DESTINAZIONE/MERCATO/TRATTA", "PROMOZIONE"])
        if key not in seen:
            cleaned.append(row)
            seen.add(key)
    return cleaned



def format_section_title(group_df: pd.DataFrame) -> str:
    """Crea il titolo della sotto-sezione: [OLTA] - [DATA MAIL] - [TIPOLOGIA]."""
    olta = cleanup_field(group_df["OLTA MITTENTE"].dropna().astype(str).iloc[0]) if not group_df.empty else "OLTA non rilevata"
    date = cleanup_field(group_df["DATA NEWSLETTER"].dropna().astype(str).iloc[0]) if not group_df.empty else "Data non rilevata"
    tipologie = [cleanup_field(x) for x in group_df["TIPOLOGIA"].dropna().astype(str).unique() if cleanup_field(x)]
    tipologia = " / ".join(tipologie) if tipologie else "Tipologia non rilevata"
    return f"{olta or 'OLTA non rilevata'} - {date or 'Data non rilevata'} - {tipologia}"


def render_confidence(confidence_value: float) -> None:
    """Mostra la confidence media della newsletter in percentuale, con colore e icona condizionale."""
    try:
        pct = round(float(confidence_value) * 100)
    except Exception:
        pct = 0

    if pct >= 90:
        color = "#16803c"
        icon = "👍"
    else:
        color = "#c62828"
        icon = "👎"

    st.markdown(
        f"<div style='font-size:0.90rem; margin-top:0.15rem;'>"
        f"<strong>Confidence media:</strong> "
        f"<span style='color:{color}; font-weight:700;'>{icon} {pct}%</span>"
        f"</div>",
        unsafe_allow_html=True,
    )


def render_source_file(file_name: str) -> None:
    """Mostra il file di origine sotto la tabella in formato piccolo."""
    st.markdown(
        f"<div style='font-size:0.78rem; color:#6b7280; margin-top:0.05rem;'>"
        f"File di origine: {file_name}"
        f"</div>",
        unsafe_allow_html=True,
    )


def normalize_logo_key(value: str) -> str:
    """Normalizza nomi OLTA/file logo per matching robusto su Streamlit Cloud."""
    value = cleanup_field(value).lower()
    value = unicodedata.normalize("NFKD", value)
    value = "".join(ch for ch in value if not unicodedata.combining(ch))
    return re.sub(r"[^a-z0-9]", "", value)


def get_logo_path(olta: str) -> Path | None:
    """Restituisce il path del logo associato all'OLTA, se disponibile.

    Cerca sia in assets/logos sia in logos e gestisce differenze di maiuscole,
    spazi, underscore e punti nei nomi file.
    """
    olta_clean = cleanup_field(olta).upper()
    expected_names = OLTA_LOGO_FILES.get(olta_clean, [])

    # 1) Tentativo diretto con le varianti note.
    for logo_dir in LOGO_DIR_CANDIDATES:
        for filename in expected_names:
            logo_path = logo_dir / filename
            if logo_path.exists():
                return logo_path

    # 2) Fallback robusto: scansione dei file immagine nelle cartelle candidate.
    target_keys = {normalize_logo_key(olta_clean)}
    target_keys.update(normalize_logo_key(Path(name).stem) for name in expected_names)

    for logo_dir in LOGO_DIR_CANDIDATES:
        if not logo_dir.exists():
            continue
        for logo_path in list(logo_dir.glob("*.png")) + list(logo_dir.glob("*.jpg")) + list(logo_dir.glob("*.jpeg")):
            if normalize_logo_key(logo_path.stem) in target_keys:
                return logo_path

    return None


def render_newsletter_header(group_df: pd.DataFrame, section_title: str) -> None:
    """Mostra logo OLTA + titolo della sotto-sezione newsletter."""
    olta = cleanup_field(group_df["OLTA MITTENTE"].dropna().astype(str).iloc[0]) if not group_df.empty else ""
    logo_path = get_logo_path(olta)

    if logo_path:
        logo_col, title_col = st.columns([1, 5])
        with logo_col:
            st.image(str(logo_path), width=LOGO_WIDTH_PX)
        with title_col:
            st.markdown(f"### {section_title}")
    else:
        st.markdown(f"### {section_title}")


def merge_edited_section(original_group: pd.DataFrame, edited_display: pd.DataFrame) -> pd.DataFrame:
    """Ricompatta le modifiche della sotto-tabella con le colonne tecniche nascoste."""
    merged = original_group.copy().reset_index(drop=True)
    edited_display = edited_display.reset_index(drop=True)

    # Se l'utente aggiunge righe nella sotto-tabella, creiamo nuove righe usando i metadati
    # della newsletter di partenza e riempiamo le colonne visibili con i valori editati.
    if len(edited_display) > len(merged):
        template = merged.iloc[0].to_dict() if not merged.empty else {col: "" for col in FINAL_COLUMNS}
        extra_rows = []
        for _ in range(len(edited_display) - len(merged)):
            new_row = template.copy()
            for col in DISPLAY_COLUMNS:
                new_row[col] = ""
            extra_rows.append(new_row)
        merged = pd.concat([merged, pd.DataFrame(extra_rows)], ignore_index=True)

    # Se l'utente elimina righe dalla sotto-tabella, manteniamo solo lo stesso numero di righe.
    if len(edited_display) < len(merged):
        merged = merged.iloc[:len(edited_display)].copy().reset_index(drop=True)

    for col in DISPLAY_COLUMNS:
        if col in edited_display.columns:
            merged[col] = edited_display[col].values

    return merged

def parse_italian_date(value: str):
    """Converte una data gg/mm/aaaa in pandas.Timestamp, se possibile."""
    try:
        return pd.to_datetime(str(value), dayfirst=True, errors="coerce")
    except Exception:
        return pd.NaT


def get_week_label(df: pd.DataFrame) -> str:
    """Restituisce una label tipo 'Week 25' sulla base delle date newsletter."""
    if df.empty or "DATA NEWSLETTER" not in df.columns:
        return "Week XX"
    parsed_dates = [parse_italian_date(x) for x in df["DATA NEWSLETTER"].dropna().astype(str)]
    parsed_dates = [x for x in parsed_dates if pd.notna(x)]
    if not parsed_dates:
        return "Week XX"
    week_numbers = [int(x.isocalendar().week) for x in parsed_dates]
    most_common_week = pd.Series(week_numbers).mode()
    if most_common_week.empty:
        return "Week XX"
    return f"Week {int(most_common_week.iloc[0])}"


def normalize_promo_for_mail(value: str) -> str:
    """Normalizza leggermente la promozione per renderla più naturale nella mail."""
    value = cleanup_field(value)
    if not value:
        return "promozione non specificata"
    # Nei testi commerciali è più leggibile 'fino al -40%' rispetto a 'fino al 40%'.
    value = re.sub(r"(?i)fino\s+al\s+(\d{1,2}\s?%)", r"fino al -\1", value)
    value = re.sub(r"(?i)fino\s+a\s+(\d{1,2}\s?%)", r"fino al -\1", value)
    value = re.sub(r"(?i)jusqu['’]à\s+(\d{1,2}\s?%)", r"fino al -\1", value)
    value = value.replace("- ", "-")
    return value


def is_gnv_row(row: pd.Series) -> bool:
    """Identifica righe in cui GNV risulta promossa in modo esplicito."""
    company = str(row.get("COMPAGNIA PROMOSSA", ""))
    subject = str(row.get("OGGETTO EMAIL", ""))
    text = f"{company} {subject}"
    return bool(re.search(r"(?i)(?<![A-Z])GNV(?![A-Z])", text))


def format_market_and_promo(row: pd.Series) -> str:
    destination = cleanup_field(row.get("DESTINAZIONE/MERCATO/TRATTA", ""))
    promo = normalize_promo_for_mail(row.get("PROMOZIONE", ""))
    if destination and promo:
        return f"{destination} ({promo})"
    if destination:
        return destination
    return promo


def format_company_items(group_df: pd.DataFrame, max_items: int = 10) -> str:
    """Crea una frase sintetica con compagnie, mercati/tratte e promozioni."""
    items = []
    for _, row in group_df.iterrows():
        company = cleanup_field(row.get("COMPAGNIA PROMOSSA", ""))
        market_promo = format_market_and_promo(row)
        if company and market_promo:
            item = f"{company} {market_promo}"
        elif company:
            item = company
        else:
            item = market_promo
        if item and item not in items:
            items.append(item)
    if not items:
        return "contenuti promozionali non chiaramente classificati"
    if len(items) > max_items:
        return "; ".join(items[:max_items]) + "; ulteriori promozioni rilevate"
    return "; ".join(items)


def create_mail_proposal(df: pd.DataFrame) -> str:
    """Genera una proposta di mail settimanale al cliente, con focus sul posizionamento GNV."""
    if df.empty:
        return "Buongiorno Davide,\nriportiamo di seguito l’analisi OLTA Newsletter settimanale per la Week XX.\n\nNessuna newsletter analizzata."

    work_df = df.copy()
    for column in FINAL_COLUMNS:
        if column not in work_df.columns:
            work_df[column] = ""

    week_label = get_week_label(work_df)
    n_files = work_df["FILE ORIGINE"].nunique()
    gnv_df = work_df[work_df.apply(is_gnv_row, axis=1)].copy()
    gnv_oltas = [x for x in gnv_df["OLTA MITTENTE"].dropna().astype(str).unique() if cleanup_field(x)]

    if len(gnv_df) > 0:
        if len(gnv_oltas) == 1:
            opening_insight = f"Settimana con presenza GNV rilevata in particolare su {gnv_oltas[0]}."
        elif len(gnv_oltas) > 1:
            opening_insight = "Settimana con presenza GNV rilevata su " + ", ".join(gnv_oltas[:-1]) + f" e {gnv_oltas[-1]}."
        else:
            opening_insight = "Settimana con presenza GNV rilevata all’interno delle newsletter analizzate."
    else:
        opening_insight = "Settimana senza evidenze GNV esplicite nelle newsletter analizzate, ma utile per monitorare il posizionamento competitivo delle altre compagnie."

    lines = [
        "Buongiorno Davide,",
        f"riportiamo di seguito l’analisi OLTA Newsletter settimanale per la {week_label}.",
        opening_insight,
        "",
    ]

    # Ordine per OLTA e data, preservando il più possibile l'ordine di caricamento/origine.
    grouped = work_df.groupby(["OLTA MITTENTE", "DATA NEWSLETTER", "TIPOLOGIA", "FILE ORIGINE"], sort=False, dropna=False)
    olta_order = [cleanup_field(x) for x in work_df["OLTA MITTENTE"].dropna().astype(str).unique() if cleanup_field(x)]

    for olta in olta_order:
        olta_df = work_df[work_df["OLTA MITTENTE"].astype(str) == olta]
        if olta_df.empty:
            continue
        gnv_present = olta_df.apply(is_gnv_row, axis=1).any()
        if gnv_present:
            lines.append(f"• {olta}: GNV presente nelle comunicazioni analizzate")
        else:
            lines.append(f"• {olta}:")

        for (g_olta, date, typology, file_origin), group in grouped:
            if cleanup_field(g_olta) != olta:
                continue
            date = cleanup_field(date) or "data non rilevata"
            group_gnv = group[group.apply(is_gnv_row, axis=1)]
            group_non_gnv = group[~group.apply(is_gnv_row, axis=1)]

            if not group_gnv.empty:
                gnv_text = format_company_items(group_gnv)
                if not group_non_gnv.empty:
                    other_text = format_company_items(group_non_gnv)
                    detail = f"GNV promossa su {gnv_text}. Il paniere include anche {other_text}."
                else:
                    detail = f"GNV promossa su {gnv_text}."
            else:
                detail = format_company_items(group)
                detail = f"comunicazione non GNV-led, con focus su {detail}."

            # Evita ripetizioni come 'GNV GNV Sicilia'.
            detail = re.sub(r"(?i)GNV promossa su GNV\s+", "GNV promossa su ", detail)
            lines.append(f"- {date}: {detail}")
        lines.append("")

    lines.append(f"Nel complesso, sono state analizzate {n_files} newsletter. La bozza è generata automaticamente con regole Python e va revisionata prima dell’invio al cliente.")
    return "\n".join(lines).strip()



def normalize_olta_for_ppt(value: str) -> str:
    """Restituisce la label OLTA coerente con il template PowerPoint."""
    key = normalize_logo_key(value)
    mapping = {
        "ferryhopper": "FERRY HOPPER",
        "aferry": "AFERRY",
        "alloferry": "ALLOFERRY",
        "directferries": "DIRECT FERRIES",
        "lacentraledesferries": "LA CENTRALE DES FERRIES",
        "lacentrale": "LA CENTRALE DES FERRIES",
        "netferry": "NETFERRY",
        "traghettiper": "TRAGHETTIPER",
        "traghetticom": "TRAGHETTI.COM",
        "adac": "ADAC",
    }
    return mapping.get(key, cleanup_field(value).upper())


def get_template_pptx_path() -> Path | None:
    """Cerca il template PowerPoint in templates/ o nella root dell'app."""
    for candidate in TEMPLATE_PPTX_CANDIDATES:
        if candidate.exists():
            return candidate
    return None


def get_week_info(df: pd.DataFrame) -> dict:
    """Calcola numero settimana e intervallo lunedì-domenica dalle date newsletter."""
    parsed_dates = []
    if not df.empty and "DATA NEWSLETTER" in df.columns:
        for value in df["DATA NEWSLETTER"].dropna().astype(str):
            dt = parse_italian_date(value)
            if pd.notna(dt):
                parsed_dates.append(dt)
    if not parsed_dates:
        return {
            "week_number": "XX",
            "week_label": "Week XX",
            "monday": None,
            "sunday": None,
            "friday": None,
            "date_range": "lunedì gg/mm/aaaa - domenica gg/mm/aaaa",
        }

    week_numbers = [int(x.isocalendar().week) for x in parsed_dates]
    week_number = int(pd.Series(week_numbers).mode().iloc[0])
    reference_dates = [x for x in parsed_dates if int(x.isocalendar().week) == week_number]
    reference = min(reference_dates) if reference_dates else min(parsed_dates)
    monday = reference - pd.Timedelta(days=int(reference.weekday()))
    sunday = monday + pd.Timedelta(days=6)
    date_range = f"lunedì {monday.strftime('%d/%m/%Y')} - domenica {sunday.strftime('%d/%m/%Y')}"
    return {
        "week_number": str(week_number),
        "week_label": f"Week {week_number}",
        "monday": monday,
        "sunday": sunday,
        "friday": sunday,
        "date_range": date_range,
    }


def get_communication_counts_for_ppt(df: pd.DataFrame) -> list[int]:
    """Conta le comunicazioni per OLTA come numero di newsletter/file, non come righe promozione."""
    counts = {olta: 0 for olta in PPT_OLTA_CATEGORIES}
    if df.empty:
        return [0 for _ in PPT_OLTA_CATEGORIES]

    if "FILE ORIGINE" in df.columns:
        newsletter_df = df.groupby("FILE ORIGINE", sort=False, dropna=False).agg({"OLTA MITTENTE": "first"}).reset_index()
    else:
        newsletter_df = df.copy()

    for value in newsletter_df.get("OLTA MITTENTE", pd.Series(dtype=str)).fillna("").astype(str):
        olta = normalize_olta_for_ppt(value)
        if olta not in counts:
            counts[olta] = 0
        counts[olta] += 1
    return [int(counts.get(olta, 0)) for olta in PPT_OLTA_CATEGORIES]


def replace_text_in_slide(slide, replacements: dict[str, str]) -> None:
    """Sostituisce testo nei textbox mantenendo la struttura semplice del template."""
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        original = shape.text
        if not original:
            continue
        updated = original
        for old, new in replacements.items():
            updated = updated.replace(old, new)
        if updated != original:
            shape.text = updated


def remove_shape(shape) -> None:
    """Rimuove una shape da una slide."""
    element = shape._element
    element.getparent().remove(element)


def duplicate_slide(prs, source_slide):
    """Duplica una slide preservando le forme del template.

    La slide dettaglio del template contiene solo forme, testo e tabella placeholder;
    non è necessario ricopiare manualmente le relazioni, che in python-pptx 1.x
    usano API interne diverse a seconda della versione.
    """
    blank_layout = prs.slide_layouts[6]
    copied_slide = prs.slides.add_slide(blank_layout)
    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        copied_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return copied_slide


def update_title_slide(slide, week_info: dict) -> None:
    replace_text_in_slide(
        slide,
        {
            "Week [numero Settimana]": f"Week {week_info['week_number']}",
            "lunedì gg/mm/aaaa - domenica gg/mm/aaaa": week_info["date_range"],
            "lunedì gg/mm/aaaa - venerdì gg/mm/aaaa": week_info["date_range"],
        },
    )


def update_summary_chart_slide(slide, week_info: dict, df: pd.DataFrame) -> None:
    replace_text_in_slide(slide, {"Week [numero Settimana]": f"Week {week_info['week_number']}"})
    counts = get_communication_counts_for_ppt(df)
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        chart = shape.chart
        chart_data = CategoryChartData()
        chart_data.categories = PPT_OLTA_CATEGORIES
        chart_data.add_series("Invii", counts)
        chart.replace_data(chart_data)
        try:
            series = chart.series[0]
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                r, g, b = PPT_BAR_COLORS[idx % len(PPT_BAR_COLORS)]
                fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception:
            pass
        break


def newsletter_group_title(group_df: pd.DataFrame) -> str:
    date = cleanup_field(group_df["DATA NEWSLETTER"].dropna().astype(str).iloc[0]) if not group_df.empty else "gg/mm/aaaa"
    tipologie = [cleanup_field(x) for x in group_df["TIPOLOGIA"].dropna().astype(str).unique() if cleanup_field(x)] if not group_df.empty else []
    tipologia = " / ".join(tipologie) if tipologie else "tipologia newsletter"
    return f"Newsletter {date} – {tipologia}"


def pack_newsletter_groups_for_ppt(df: pd.DataFrame) -> list[list[pd.DataFrame]]:
    """Raggruppa le newsletter in slide. Max 2 tabelle per slide se stessa OLTA e tabelle brevi."""
    groups = []
    for _, group_df in df.groupby("FILE ORIGINE", sort=False, dropna=False):
        groups.append(group_df.reset_index(drop=True))

    chunks = []
    i = 0
    while i < len(groups):
        current = groups[i]
        current_olta = cleanup_field(current["OLTA MITTENTE"].dropna().astype(str).iloc[0]) if not current.empty else ""
        if i + 1 < len(groups):
            nxt = groups[i + 1]
            next_olta = cleanup_field(nxt["OLTA MITTENTE"].dropna().astype(str).iloc[0]) if not nxt.empty else ""
            same_olta = normalize_olta_for_ppt(current_olta) == normalize_olta_for_ppt(next_olta)
            if same_olta and len(current) <= 5 and len(nxt) <= 5:
                chunks.append([current, nxt])
                i += 2
                continue
        chunks.append([current])
        i += 1
    return chunks


def set_cell_text(cell, value: str, bold: bool = False, font_size: int = 10, header: bool = False) -> None:
    value = cleanup_field(value)
    cell.text = value
    try:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER if header else PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.bold = bold
            run.font.size = Pt(font_size)
            run.font.name = "Arial"
            if header:
                run.font.color.rgb = RGBColor(255, 255, 255)
            else:
                run.font.color.rgb = RGBColor(35, 35, 35)


def add_ppt_table(slide, group_df: pd.DataFrame, left, top, width, height, compact: bool = False) -> None:
    """Aggiunge una tabella PowerPoint con colonne Compagnia / Destinazione / Sconto."""
    rows = max(2, len(group_df) + 1)
    cols = 3
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = int(width * 0.25)
    table.columns[1].width = int(width * 0.50)
    table.columns[2].width = int(width * 0.25)

    header_font = 9 if compact else 10
    body_font = 8 if compact else 10
    headers = ["COMPAGNIA", "DESTINAZIONE / TRATTA / MERCATO", "SCONTO"]
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
        set_cell_text(cell, header, bold=True, font_size=header_font, header=True)

    for r_idx, (_, row) in enumerate(group_df.iterrows(), start=1):
        values = [
            row.get("COMPAGNIA PROMOSSA", ""),
            row.get("DESTINAZIONE/MERCATO/TRATTA", ""),
            row.get("PROMOZIONE", ""),
        ]
        for c, value in enumerate(values):
            cell = table.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            set_cell_text(cell, str(value), bold=False, font_size=body_font, header=False)


def find_text_shape(slide, pattern: str):
    regex = re.compile(pattern, re.I)
    for shape in slide.shapes:
        if hasattr(shape, "text") and regex.search(shape.text or ""):
            return shape
    return None


def add_subtitle_text(slide, text: str, left, top, width, height) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.color.rgb = RGBColor(60, 60, 60)


def replace_logo_placeholder(slide, olta: str) -> None:
    placeholder = find_text_shape(slide, r"\[LOGO OLTA\]")
    left = placeholder.left if placeholder is not None else Inches(1.33)
    top = placeholder.top if placeholder is not None else Inches(0.55)
    width = placeholder.width if placeholder is not None else Inches(4.0)
    height = placeholder.height if placeholder is not None else Inches(0.55)
    if placeholder is not None:
        remove_shape(placeholder)

    logo_path = get_logo_path(olta)
    if logo_path:
        try:
            slide.shapes.add_picture(str(logo_path), left, top, height=height)
            return
        except Exception:
            pass
    add_subtitle_text(slide, normalize_olta_for_ppt(olta), left, top, width, height)


def populate_detail_slide(slide, newsletter_groups: list[pd.DataFrame]) -> None:
    """Popola una slide dettaglio con una o due newsletter della stessa OLTA."""
    if not newsletter_groups:
        return
    first_group = newsletter_groups[0]
    olta = cleanup_field(first_group["OLTA MITTENTE"].dropna().astype(str).iloc[0]) if not first_group.empty else ""
    replace_logo_placeholder(slide, olta)

    for shape in list(slide.shapes):
        if getattr(shape, "has_table", False):
            remove_shape(shape)

    subtitle_shape = find_text_shape(slide, r"Newsletter .*\[tipologia newsletter\]|Newsletter gg/mm/aaaa")
    if subtitle_shape is not None:
        subtitle_shape.text = newsletter_group_title(first_group)
        subtitle_left, subtitle_width, subtitle_height = subtitle_shape.left, subtitle_shape.width, subtitle_shape.height
    else:
        subtitle_left, subtitle_width, subtitle_height = Inches(1.33), Inches(12.5), Inches(0.6)
        add_subtitle_text(slide, newsletter_group_title(first_group), subtitle_left, Inches(1.37), subtitle_width, subtitle_height)

    table_left = Inches(1.33)
    table_width = Inches(17.35)

    if len(newsletter_groups) == 1:
        add_ppt_table(slide, first_group, table_left, Inches(2.55), table_width, Inches(6.75), compact=len(first_group) > 7)
    else:
        add_ppt_table(slide, first_group, table_left, Inches(2.15), table_width, Inches(2.75), compact=True)
        second_group = newsletter_groups[1]
        add_subtitle_text(slide, newsletter_group_title(second_group), subtitle_left, Inches(5.20), subtitle_width, subtitle_height)
        add_ppt_table(slide, second_group, table_left, Inches(5.95), table_width, Inches(3.00), compact=True)


def build_powerpoint(df: pd.DataFrame, template_bytes: bytes | None = None) -> BytesIO:
    """Genera il PowerPoint OLTA sulla base del template prestabilito."""
    if Presentation is None:
        raise RuntimeError("La libreria python-pptx non è installata. Aggiungi python-pptx al requirements.txt.")

    if template_bytes:
        prs = Presentation(BytesIO(template_bytes))
    else:
        template_path = get_template_pptx_path()
        if template_path is None:
            raise FileNotFoundError("Template PowerPoint non trovato. Carica templates/template_olta.pptx nella repository o usa l'uploader template nell'app.")
        prs = Presentation(str(template_path))

    if len(prs.slides) < 3:
        raise ValueError("Il template PowerPoint deve contenere almeno 3 slide: cover, grafico riepilogo e slide dettaglio newsletter.")

    week_info = get_week_info(df)
    update_title_slide(prs.slides[0], week_info)
    update_summary_chart_slide(prs.slides[1], week_info, df)

    chunks = pack_newsletter_groups_for_ppt(df)
    detail_template_slide = prs.slides[2]
    detail_slides = [detail_template_slide]
    for _ in range(max(0, len(chunks) - 1)):
        detail_slides.append(duplicate_slide(prs, detail_template_slide))

    if chunks:
        for slide, chunk in zip(detail_slides, chunks):
            populate_detail_slide(slide, chunk)
    else:
        populate_detail_slide(detail_template_slide, [])

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def build_excel(df: pd.DataFrame, mail_text: str) -> BytesIO:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        export_df = df.copy()
        for column in EXPORT_COLUMNS:
            if column not in export_df.columns:
                export_df[column] = ""
        export_df = export_df[EXPORT_COLUMNS]
        export_df.to_excel(writer, sheet_name="Database", index=False)
        summary_rows = []
        if not df.empty:
            summary_rows.append({"Indicatore": "Newsletter caricate", "Valore": df["FILE ORIGINE"].nunique()})
            summary_rows.append({"Indicatore": "Righe estratte", "Valore": len(df)})
            for olta, count in df["OLTA MITTENTE"].value_counts().items():
                summary_rows.append({"Indicatore": f"Comunicazioni - {olta}", "Valore": count})
            for typ, count in df["TIPOLOGIA"].value_counts().items():
                summary_rows.append({"Indicatore": f"Tipologia - {typ}", "Valore": count})
        summary_df = pd.DataFrame(summary_rows or [{"Indicatore": "Note", "Valore": "Nessun dato"}])
        summary_df.to_excel(writer, sheet_name="Summary", index=False)
        pd.DataFrame([{"Proposta Mail": mail_text}]).to_excel(writer, sheet_name="Proposta Mail", index=False)
    output.seek(0)
    return output


st.set_page_config(page_title="OLTA Newsletter Extractor", layout="wide")
st.title("OLTA Newsletter Extractor")
st.caption("Versione 0.6 — estrazione rule-based con output Excel e PowerPoint da template")

st.markdown(
    "Carica le newsletter in formato `.msg`, `.eml`, `.html` o `.txt`. "
    "L'app creerà una sotto-sezione per ciascuna newsletter, con tabella modificabile, confidence, file di origine, proposta mail e PowerPoint finale da template."
)

uploaded_files = st.file_uploader(
    "Carica newsletter",
    type=["msg", "eml", "html", "htm", "txt"],
    accept_multiple_files=True,
)

with st.expander("Template PowerPoint", expanded=False):
    st.caption(
        "Di default l'app usa `templates/template_olta.pptx`. "
        "Se il template non è presente nella repository, puoi caricarlo manualmente qui."
    )
    uploaded_template_pptx = st.file_uploader(
        "Carica template PowerPoint opzionale",
        type=["pptx"],
        accept_multiple_files=False,
        key="template_pptx_uploader",
    )

with st.expander("Diagnostica loghi e template", expanded=False):
    st.caption("Usa questa sezione se i loghi o il template PowerPoint non vengono letti su Streamlit Cloud.")
    checked_dirs = []
    for logo_dir in LOGO_DIR_CANDIDATES:
        checked_dirs.append({
            "Percorso controllato": str(logo_dir),
            "Esiste": logo_dir.exists(),
            "File trovati": ", ".join(sorted([p.name for p in logo_dir.glob("*")])) if logo_dir.exists() else "",
        })
    st.dataframe(pd.DataFrame(checked_dirs), use_container_width=True, hide_index=True)

    checked_templates = []
    for template_path in TEMPLATE_PPTX_CANDIDATES:
        checked_templates.append({
            "Template controllato": str(template_path),
            "Esiste": template_path.exists(),
        })
    st.dataframe(pd.DataFrame(checked_templates), use_container_width=True, hide_index=True)

if uploaded_files:
    all_rows = []
    errors = []
    for uploaded_file in uploaded_files:
        try:
            parsed = parse_uploaded_file(uploaded_file)
            all_rows.extend(extract_rows(parsed, uploaded_file.name))
        except Exception as exc:
            errors.append(f"{uploaded_file.name}: {exc}")

    if errors:
        st.warning("Alcuni file non sono stati letti correttamente:\n" + "\n".join(errors))

    df = pd.DataFrame(all_rows, columns=FINAL_COLUMNS)
    if df.empty:
        st.info("Nessuna informazione estratta dai file caricati.")
    else:
        st.subheader("Newsletter estratte - tabelle modificabili")
        st.write(
            "Ogni sezione corrisponde a una singola newsletter. "
            "Le colonne tecniche non sono mostrate nella tabella, ma vengono usate per il riepilogo e l'export."
        )

        edited_groups = []
        grouped = df.groupby("FILE ORIGINE", sort=False, dropna=False)

        for section_idx, (file_origin, group_df) in enumerate(grouped, start=1):
            group_df = group_df.reset_index(drop=True)
            section_title = format_section_title(group_df)

            render_newsletter_header(group_df, section_title)

            display_df = group_df.copy()
            for column in DISPLAY_COLUMNS:
                if column not in display_df.columns:
                    display_df[column] = ""
            display_df = display_df[DISPLAY_COLUMNS]

            edited_display = st.data_editor(
                display_df,
                use_container_width=True,
                num_rows="dynamic",
                key=f"newsletter_table_{section_idx}_{file_origin}",
                column_config={
                    "COMPAGNIA PROMOSSA": st.column_config.TextColumn("COMPAGNIA PROMOSSA"),
                    "DESTINAZIONE/MERCATO/TRATTA": st.column_config.TextColumn("DESTINAZIONE/MERCATO/TRATTA"),
                    "PROMOZIONE": st.column_config.TextColumn("PROMOZIONE"),
                    "OGGETTO EMAIL": st.column_config.TextColumn("OGGETTO EMAIL"),
                },
            )

            confidence_mean = group_df["CONFIDENCE"].dropna().astype(float).mean() if "CONFIDENCE" in group_df else 0
            render_confidence(confidence_mean)
            render_source_file(str(file_origin))

            edited_group = merge_edited_section(group_df, edited_display)
            edited_groups.append(edited_group)

            st.divider()

        edited_df = pd.concat(edited_groups, ignore_index=True) if edited_groups else pd.DataFrame(columns=FINAL_COLUMNS)

        mail_text = create_mail_proposal(edited_df)
        st.subheader("Proposta Mail")
        mail_text = st.text_area("Bozza mail al cliente", value=mail_text, height=360)

        col_download_excel, col_download_ppt = st.columns(2)

        excel_data = build_excel(edited_df, mail_text)
        with col_download_excel:
            st.download_button(
                "Scarica Excel",
                data=excel_data,
                file_name="OLTA_newsletter_output.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

        with col_download_ppt:
            try:
                template_bytes = uploaded_template_pptx.getvalue() if uploaded_template_pptx is not None else None
                pptx_data = build_powerpoint(edited_df, template_bytes=template_bytes)
                week_label_for_file = get_week_info(edited_df)["week_label"].replace(" ", "_")
                st.download_button(
                    "Scarica PowerPoint",
                    data=pptx_data,
                    file_name=f"OLTA_newsletter_{week_label_for_file}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as exc:
                st.error(f"PowerPoint non generato: {exc}")
else:
    st.info("Carica uno o più file per iniziare.")
